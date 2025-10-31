import os
import sys
import json
import logging
import tempfile
import shutil
from pathlib import Path
from mcp.server.fastmcp import FastMCP
from pptx import Presentation
from pptx.util import Inches, Pt
from playwright.async_api import async_playwright
from dotenv import load_dotenv
from typing import List, Dict, Any # Import typing helpers

# Configure logging
logging.basicConfig(level=logging.INFO, stream=sys.stderr, format='%(asctime)s - %(levelname)s - %(message)s')

# Load environment variables (though none are strictly needed for this server)
load_dotenv()

# --- Server Setup ---
mcp = FastMCP(
    name="presentation-creator-server",
    instructions="Assembles slide content into a PowerPoint presentation.",
)

# --- Constants ---
# Define standard slide dimensions (16:9 aspect ratio in inches)
PPTX_WIDTH = Inches(10)
PPTX_HEIGHT = Inches(5.625)

# --- Tool Implementation: Assemble Presentation ---
@mcp.tool()
async def assemble_presentation(slides_html: List[str], filename: str = "presentation") -> str: # Changed input name and type
    """
    Assembles a PowerPoint presentation from a list of HTML strings.
    Saves each HTML string to a file, screenshots it, and adds the image to a PPTX file.

    Args:
        slides_html: A list where each item is a string containing the full HTML code for one slide.
        filename: The desired base name for the output PPTX file (without extension).

    Returns:
        The absolute path to the generated .pptx file or an error message.
    """
    logging.info(f"Received request to assemble presentation: {filename}.pptx with {len(slides_html)} slides.")
    temp_dir = None
    try:
        if not isinstance(slides_html, list):
             return "Error: Input slides_html must be a list."
        if not all(isinstance(item, str) for item in slides_html):
             return "Error: All items in slides_html must be strings."

        if not slides_html:
            return "Error: No slide HTML data provided."

        # Create a temporary directory for HTML files and screenshots
        temp_dir = tempfile.mkdtemp(prefix="mcp_ppt_")
        logging.info(f"Created temporary directory: {temp_dir}")
        image_paths = []

        # Launch Playwright
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()
            await page.set_viewport_size({"width": 800, "height": 450})

            # Process each slide
            for i, html_content in enumerate(slides_html): # Iterate through provided HTML strings
                slide_num = i + 1
                html_path = os.path.join(temp_dir, f"slide_{slide_num}.html")
                img_path = os.path.join(temp_dir, f"slide_{slide_num}.png")

                # Write the provided HTML content to file
                with open(html_path, "w", encoding="utf-8") as f:
                    f.write(html_content)

                await page.goto(f"file://{html_path}")
                await page.screenshot(path=img_path)
                image_paths.append(img_path)
                logging.info(f"Generated screenshot for slide {slide_num}: {img_path}")

            await browser.close()

        # Create PowerPoint presentation
        prs = Presentation()
        prs.slide_width = PPTX_WIDTH
        prs.slide_height = PPTX_HEIGHT
        blank_slide_layout = prs.slide_layouts[6]

        for img_path in image_paths:
            slide = prs.slides.add_slide(blank_slide_layout)
            left = top = Inches(0)
            pic = slide.shapes.add_picture(img_path, left, top, width=prs.slide_width, height=prs.slide_height)

        # Save presentation to Downloads folder
        downloads_path = Path.home() / "Downloads"
        downloads_path.mkdir(parents=True, exist_ok=True)
        output_filename = f"{filename}.pptx"
        output_path = downloads_path / output_filename
        prs.save(output_path)
        final_path = str(output_path.resolve())
        logging.info(f"Presentation saved to: {final_path}")

        return f"Presentation successfully created and saved to: {final_path}"

    except Exception as e:
        logging.exception("An error occurred during presentation assembly.")
        return f"Error assembling presentation: {e}"
    finally:
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
                logging.info(f"Cleaned up temporary directory: {temp_dir}")
            except Exception as e:
                logging.error(f"Error cleaning up temporary directory {temp_dir}: {e}")


# --- PDF Creation Tool ---
@mcp.tool()
async def create_pdf_from_html(html_content: str, filename: str = "document") -> str:
    """
    Generates a PDF document from a string containing HTML code.

    Args:
        html_content: A string containing the full HTML code for the document.
        filename: The desired base name for the output PDF file (without extension).

    Returns:
        The absolute path to the generated .pdf file or an error message.
    """
    logging.info(f"Received request to create PDF: {filename}.pdf")
    temp_dir = None
    html_path = None
    try:
        if not isinstance(html_content, str) or not html_content.strip():
            return "Error: html_content must be a non-empty string."

        # Create a temporary directory and file for the HTML
        temp_dir = tempfile.mkdtemp(prefix="mcp_pdf_")
        html_path = os.path.join(temp_dir, "content.html")
        with open(html_path, "w", encoding="utf-8") as f:
            f.write(html_content)
        logging.info(f"Saved HTML to temporary file: {html_path}")

        # Launch Playwright and generate PDF
        async with async_playwright() as p:
            browser = await p.chromium.launch()
            page = await browser.new_page()

            await page.goto(f"file://{html_path}")

            # Generate PDF
            downloads_path = Path.home() / "Downloads"
            downloads_path.mkdir(parents=True, exist_ok=True)
            output_filename = f"{filename}.pdf"
            output_path = downloads_path / output_filename
            pdf_options = {
                "path": str(output_path),
                "format": "A4", # Standard paper size
                "printBackground": True, # Include background colors/images
                "margin": { # Optional margins
                    "top": "20mm",
                    "bottom": "20mm",
                    "left": "20mm",
                    "right": "20mm"
                }
            }
            await page.pdf(**pdf_options)
            await browser.close()

        final_path = str(output_path.resolve())
        logging.info(f"PDF saved to: {final_path}")
        return f"PDF successfully created and saved to: {final_path}"

    except Exception as e:
        logging.exception("An error occurred during PDF creation.")
        return f"Error creating PDF: {e}"
    finally:
        # Clean up temporary directory/file
        if temp_dir and os.path.exists(temp_dir):
            try:
                shutil.rmtree(temp_dir)
                logging.info(f"Cleaned up temporary directory: {temp_dir}")
            except Exception as e:
                logging.error(f"Error cleaning up temporary directory {temp_dir}: {e}")


# --- Server Execution ---
if __name__ == "__main__":
    logging.info("Starting Presentation Creator Server...")
    mcp.run(transport='stdio')
    logging.info("Presentation Creator Server stopped.")
